"""
Extract PPTX content to Markdown + images.

Usage:
    python utils/extract_pptx.py <input.pptx> [output_dir]

Output structure:
    output_dir/
        index.md          — full markdown with slide content
        images/
            slide_01_img_0.png
            ...
"""

import argparse
import re
import sys
from pathlib import Path

# Fix Windows console encoding for Cyrillic paths
if sys.stdout.encoding and sys.stdout.encoding.lower() not in ("utf-8", "utf-8-sig"):
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from PIL import Image
import io


def slug(text: str) -> str:
    text = re.sub(r"[^\w\s-]", "", text, flags=re.UNICODE)
    return re.sub(r"\s+", "_", text).strip("_").lower()


def extract_pptx(pptx_path: Path, output_dir: Path) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    images_dir = output_dir / "images"
    images_dir.mkdir(exist_ok=True)

    prs = Presentation(str(pptx_path))
    md_lines: list[str] = [f"# {pptx_path.stem}\n"]

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_label = f"Слайд {slide_idx}"

        # Try to get slide title
        title_text = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            slide_label = title_text

        md_lines.append(f"\n## Слайд {slide_idx}: {title_text or '(без заголовка)'}\n")

        img_counter = 0
        text_blocks: list[str] = []

        for shape in slide.shapes:
            # ── Text ──────────────────────────────────────────────────────────
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if not line:
                        continue
                    # Detect heading by font size or if it's the title shape
                    is_title = shape == slide.shapes.title
                    font_size = None
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size.pt
                            break

                    if is_title:
                        continue  # already used as slide heading
                    elif font_size and font_size >= 24:
                        text_blocks.append(f"### {line}")
                    elif para.level == 0:
                        text_blocks.append(line)
                    else:
                        indent = "  " * para.level
                        text_blocks.append(f"{indent}- {line}")

            # ── Images ────────────────────────────────────────────────────────
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_data = shape.image.blob
                ext = shape.image.ext  # e.g. "jpeg", "png"
                img_name = f"slide_{slide_idx:02d}_img_{img_counter}.{ext}"
                img_path = images_dir / img_name
                img_path.write_bytes(img_data)
                rel_path = f"images/{img_name}"
                text_blocks.append(f"![{img_name}]({rel_path})")
                img_counter += 1

            # ── Tables ────────────────────────────────────────────────────────
            if shape.has_table:
                tbl = shape.table
                rows_md: list[str] = []
                for row_idx, row in enumerate(tbl.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows_md.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        rows_md.append("|" + "|".join(["---"] * len(cells)) + "|")
                text_blocks.extend(rows_md)

        md_lines.extend(text_blocks)

    md_path = output_dir / "index.md"
    md_path.write_text("\n".join(md_lines), encoding="utf-8")
    print(f"[pptx] Done: {md_path}  (images: {images_dir})")


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract PPTX → Markdown + images")
    parser.add_argument("pptx", help="Path to .pptx file")
    parser.add_argument("output", nargs="?", help="Output directory (default: <pptx_stem>/)")
    args = parser.parse_args()

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    output_dir = Path(args.output) if args.output else pptx_path.parent / pptx_path.stem
    extract_pptx(pptx_path, output_dir)


if __name__ == "__main__":
    main()
